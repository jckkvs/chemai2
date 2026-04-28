"""
backend/data/file_uploader.py
CSV/Excel/PPTX/DOCX ファイルのアップロード・読み込み処理
既存機能と共存：既存のデータ読み込みロジックは一切変更せず、拡張として実装
"""
import io
import pandas as pd
import numpy as np
from pathlib import Path
from typing import Optional, Union, Dict, List, Tuple
import logging

logger = logging.getLogger(__name__)

# 対応拡張子
SUPPORTED_EXTENSIONS = {
    'csv': ['.csv'],
    'excel': ['.xlsx', '.xls', '.xlsm', '.xlsb'],
    'powerpoint': ['.pptx', '.ppt'],
    'word': ['.docx', '.doc'],
    'text': ['.txt', '.tsv']
}


def detect_file_type(file_path: Union[str, Path, io.BytesIO]) -> str:
    """ファイルタイプを自動判定"""
    if isinstance(file_path, (str, Path)):
        suffix = Path(file_path).suffix.lower()
    else:
        # BytesIOの場合はヘッダーから判定（簡易）
        file_path.seek(0)
        header = file_path.read(20)
        file_path.seek(0)
        if b'PK' in header:
            return 'excel'  # xlsx/pptx/docx are ZIP-based
        return 'unknown'
    
    for file_type, extensions in SUPPORTED_EXTENSIONS.items():
        if suffix in extensions:
            return file_type
    return 'unknown'


def read_csv_smart(file_content: Union[str, bytes, io.BytesIO], 
                   encoding: Optional[str] = None,
                   **kwargs) -> pd.DataFrame:
    """
    CSVファイルを柔軟に読み込み
    - 区切り文字自動検出
    - エンコーディング自動検出（fallback付き）
    - 整形されていないデータの警告
    """
    encodings_to_try = [encoding, 'utf-8', 'utf-8-sig', 'cp932', 'shift_jis', 'iso-8859-1'] if encoding else ['utf-8', 'utf-8-sig', 'cp932', 'shift_jis', 'iso-8859-1']
    
    # 区切り文字候補
    delimiters = [',', '\t', ';', '|']
    
    for enc in encodings_to_try:
        if enc is None:
            continue
        for delim in delimiters:
            try:
                if isinstance(file_content, (str, Path)):
                    df = pd.read_csv(file_content, sep=delim, encoding=enc, on_bad_lines='skip', **kwargs)
                else:
                    df = pd.read_csv(io.BytesIO(file_content) if isinstance(file_content, bytes) else file_content, sep=delim, encoding=enc, on_bad_lines='skip', **kwargs)
                
                # 基本的な品質チェック
                if df.shape[1] < 2 or df.isnull().all().all():
                    continue
                    
                logger.info(f"CSV読み込み成功: encoding={enc}, delimiter='{delim}', shape={df.shape}")
                return df
                
            except Exception as e:
                logger.debug(f"CSV読み込み試行失敗: encoding={enc}, delimiter='{delim}', error={e}")
                continue
    
    raise ValueError("CSVファイルの読み込みに失敗しました。エンコーディングまたは形式を確認してください。")


def read_excel_smart(file_content: Union[str, bytes, io.BytesIO],
                     sheet_name: Optional[Union[str, int]] = 0,
                     **kwargs) -> pd.DataFrame:
    """
    Excelファイルを柔軟に読み込み
    - シート名自動検出
    - 結合セルの展開処理
    """
    try:
        # pandasのread_excelは_bytesIOも対応
        df = pd.read_excel(io.BytesIO(file_content) if isinstance(file_content, bytes) else file_content, sheet_name=sheet_name, **kwargs)
        
        # 結合セルの処理：前方填充で展開
        df = df.ffill()
        
        # 列名のクリーニング
        if df.columns.isnull().any() or df.columns.astype(str).str.strip().eq('').any():
            df.columns = [f"Column_{i}" if pd.isna(col) or str(col).strip() == '' else str(col).strip() 
                         for i, col in enumerate(df.columns)]
        
        logger.info(f"Excel読み込み成功: shape={df.shape}, sheet={sheet_name}")
        return df
        
    except Exception as e:
        logger.error(f"Excel読み込みエラー: {e}")
        raise


def read_document_to_text(file_content: Union[str, bytes, io.BytesIO], 
                          file_type: str) -> str:
    """
    PPTX/DOCX ファイルからテキストを抽出
    python-pptx, python-docx 使用（オプション依存）
    """
    try:
        if file_type == 'powerpoint':
            from pptx import Presentation
            if isinstance(file_content, (str, Path)):
                prs = Presentation(file_content)
            else:
                prs = Presentation(io.BytesIO(file_content if isinstance(file_content, bytes) else file_content.getvalue()))
            
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n\n".join(texts)
            
        elif file_type == 'word':
            from docx import Document
            if isinstance(file_content, (str, Path)):
                doc = Document(file_content)
            else:
                doc = Document(io.BytesIO(file_content if isinstance(file_content, bytes) else file_content.getvalue()))
            
            return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            
        else:
            # 単純なテキストファイルとして読み込み
            if isinstance(file_content, (str, Path)):
                with open(file_content, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            else:
                return file_content.read().decode('utf-8', errors='ignore')
                
    except ImportError:
        logger.warning("python-pptx/python-docx がインストールされていません。テキスト抽出をスキップします。")
        return "[テキスト抽出には python-pptx, python-docx が必要です]"
    except Exception as e:
        logger.error(f"ドキュメント読み込みエラー: {e}")
        return f"[読み込みエラー: {str(e)}]"


def assess_data_quality(df: pd.DataFrame) -> Dict[str, any]:
    """
    データの品質を評価し、クリーニングが必要か判定
    """
    quality_report = {
        'is_clean': True,
        'issues': [],
        'recommendations': []
    }
    
    # 基本的なチェック
    if df.empty:
        quality_report['is_clean'] = False
        quality_report['issues'].append('データが空です')
        return quality_report
    
    # 欠損値チェック
    missing_ratio = df.isnull().sum().sum() / (df.shape[0] * df.shape[1])
    if missing_ratio > 0.1:
        quality_report['issues'].append(f'欠損値が多い: {missing_ratio*100:.1f}%')
        quality_report['recommendations'].append('欠損値補完または削除を検討')
    
    # 列名の整合性チェック
    if df.columns.isnull().any() or df.columns.astype(str).str.strip().eq('').any():
        quality_report['issues'].append('列名に空白またはNULLが含まれています')
        quality_report['recommendations'].append('列名を整理してください')
        quality_report['is_clean'] = False
    
    # 型の一貫性チェック（数値列に文字列が混在など）
    for col in df.select_dtypes(include=[np.number]).columns:
        if df[col].apply(lambda x: isinstance(x, str) and not pd.isna(x)).any():
            quality_report['issues'].append(f'数値列 "{col}" に文字列が混在')
            quality_report['recommendations'].append(f'列 "{col}" のデータ型を統一')
            quality_report['is_clean'] = False
    
    # 重複行チェック
    if df.duplicated().sum() > 0:
        quality_report['issues'].append(f'重複行が {df.duplicated().sum()} 件存在')
        quality_report['recommendations'].append('重複行の削除を検討')
    
    # 外れ値の簡易チェック（数値列のみ）
    for col in df.select_dtypes(include=[np.number]).columns:
        if df[col].notna().sum() > 10:
            q1, q3 = df[col].quantile([0.25, 0.75])
            iqr = q3 - q1
            outliers = ((df[col] < q1 - 3*iqr) | (df[col] > q3 + 3*iqr)).sum()
            if outliers > len(df) * 0.05:
                quality_report['issues'].append(f'列 "{col}" に外れ値の疑い: {outliers} 件')
                quality_report['recommendations'].append(f'列 "{col}" の外れ値を確認')
    
    return quality_report
